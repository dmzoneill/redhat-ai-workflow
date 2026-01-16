#!/usr/bin/env python3
"""
Generate PowerPoint slides from markdown outline.

Usage:
    python scripts/generate_slides.py docs/slides/onboarding-outline.md
    python scripts/generate_slides.py docs/slides/onboarding-outline.md -o presentation.pptx

Requirements:
    pip install python-pptx
"""

import argparse
import re
from pathlib import Path
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def RgbColor(r: int, g: int, b: int):
    """Create an RGB color from components."""
    from pptx.dml.color import _Color
    # In python-pptx, we set colors via the .rgb property with an RGBColor
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


@dataclass
class Slide:
    """Represents a single slide."""
    title: str
    subtitle: str = ""
    content: list = field(default_factory=list)
    speaker_notes: str = ""
    slide_type: str = "content"  # title, content, code, diagram


def parse_markdown(md_content: str) -> list[Slide]:
    """Parse markdown outline into slide objects."""
    slides = []
    current_slide = None
    in_code_block = False
    code_content = []
    in_speaker_notes = False
    
    lines = md_content.split('\n')
    i = 0
    
    while i < len(lines):
        line = lines[i]
        
        # Skip front matter and meta sections
        if line.startswith('# AI Workflow Engineer Onboarding Slides'):
            i += 1
            continue
        if line.startswith('> **Purpose**') or line.startswith('> **Duration**') or line.startswith('> **Audience**'):
            i += 1
            continue
            
        # New slide marker: ### Slide N: Title
        slide_match = re.match(r'^### Slide \d+: (.+)$', line)
        if slide_match:
            if current_slide:
                slides.append(current_slide)
            current_slide = Slide(title=slide_match.group(1).strip())
            in_speaker_notes = False
            i += 1
            continue
        
        # Part header - create section slide
        part_match = re.match(r'^## Part \d+: (.+)$', line)
        if part_match:
            if current_slide:
                slides.append(current_slide)
            current_slide = Slide(
                title=part_match.group(1).strip(),
                slide_type="section"
            )
            i += 1
            continue
        
        # Skip appendix and tips sections
        if line.startswith('## Appendix') or line.startswith('## Tips for Google Slides'):
            break
            
        # Skip horizontal rules
        if line.strip() == '---':
            i += 1
            continue
        
        if current_slide is None:
            i += 1
            continue
            
        # Code blocks
        if line.startswith('```'):
            if in_code_block:
                # End of code block
                if code_content:
                    current_slide.content.append({
                        'type': 'code',
                        'content': '\n'.join(code_content)
                    })
                code_content = []
                in_code_block = False
            else:
                # Start of code block
                in_code_block = True
            i += 1
            continue
            
        if in_code_block:
            code_content.append(line)
            i += 1
            continue
        
        # Speaker notes
        if line.startswith('> **Speaker Notes**:') or line.startswith('> Speaker Notes:'):
            in_speaker_notes = True
            notes_text = line.split(':', 1)[1].strip() if ':' in line else ""
            current_slide.speaker_notes = notes_text
            i += 1
            continue
            
        if in_speaker_notes and line.startswith('>'):
            current_slide.speaker_notes += " " + line[1:].strip()
            i += 1
            continue
        elif in_speaker_notes and not line.startswith('>'):
            in_speaker_notes = False
        
        # Title/Subtitle
        if line.startswith('**Title**:'):
            # Already have title from slide header
            i += 1
            continue
            
        if line.startswith('**Subtitle**:'):
            current_slide.subtitle = line.replace('**Subtitle**:', '').strip()
            i += 1
            continue
        
        # Key points / bullet lists
        if line.startswith('- ') or line.startswith('* '):
            bullet_text = line[2:].strip()
            # Clean up markdown formatting
            bullet_text = re.sub(r'\*\*(.+?)\*\*', r'\1', bullet_text)  # Remove bold
            bullet_text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', bullet_text)  # Remove links
            current_slide.content.append({
                'type': 'bullet',
                'content': bullet_text
            })
            i += 1
            continue
        
        # Tables - convert to text
        if line.startswith('|') and '|' in line[1:]:
            # Parse table
            table_rows = []
            while i < len(lines) and lines[i].startswith('|'):
                row = lines[i]
                if not row.startswith('|--') and not row.startswith('| --'):
                    cells = [c.strip() for c in row.split('|')[1:-1]]
                    if cells:
                        table_rows.append(cells)
                i += 1
            if table_rows:
                current_slide.content.append({
                    'type': 'table',
                    'content': table_rows
                })
            continue
        
        # Regular text paragraphs (bold headers)
        if line.startswith('**') and line.endswith('**:'):
            header = line.strip('*:').strip()
            current_slide.content.append({
                'type': 'header',
                'content': header
            })
            i += 1
            continue
            
        if line.startswith('**') and '**:' in line:
            # Header with inline content
            parts = line.split('**:', 1)
            if len(parts) == 2:
                header = parts[0].strip('*').strip()
                content = parts[1].strip()
                current_slide.content.append({
                    'type': 'header',
                    'content': f"{header}: {content}"
                })
            i += 1
            continue
        
        i += 1
    
    # Don't forget the last slide
    if current_slide:
        slides.append(current_slide)
    
    return slides


def create_presentation(slides: list[Slide], output_path: Path) -> None:
    """Create PowerPoint presentation from slides."""
    prs = Presentation()
    
    # Set slide dimensions (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Color scheme (purple/indigo theme)
    TITLE_COLOR = RgbColor(99, 102, 241)      # Indigo-500
    ACCENT_COLOR = RgbColor(139, 92, 246)     # Purple-500
    TEXT_COLOR = RgbColor(31, 41, 55)         # Gray-800
    LIGHT_BG = RgbColor(243, 244, 246)        # Gray-100
    CODE_BG = RgbColor(249, 250, 251)         # Gray-50
    
    for slide_data in slides:
        if slide_data.slide_type == "section":
            # Section divider slide
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)
            
            # Add background shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                prs.slide_width, prs.slide_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = TITLE_COLOR
            shape.line.fill.background()
            
            # Add title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.5),
                Inches(12.333), Inches(2)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.size = Pt(54)
            p.font.bold = True
            p.font.color.rgb = RgbColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
        else:
            # Content slide
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3),
                Inches(12.333), Inches(0.8)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = TITLE_COLOR
            
            # Subtitle if present
            y_offset = Inches(1.2)
            if slide_data.subtitle:
                sub_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.1),
                    Inches(12.333), Inches(0.5)
                )
                tf = sub_box.text_frame
                p = tf.paragraphs[0]
                p.text = slide_data.subtitle
                p.font.size = Pt(20)
                p.font.color.rgb = TEXT_COLOR
                y_offset = Inches(1.7)
            
            # Content
            content_box = slide.shapes.add_textbox(
                Inches(0.5), y_offset,
                Inches(12.333), Inches(5.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            first_para = True
            for item in slide_data.content:
                if item['type'] == 'bullet':
                    if first_para:
                        p = tf.paragraphs[0]
                        first_para = False
                    else:
                        p = tf.add_paragraph()
                    p.text = f"‚Ä¢ {item['content']}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = TEXT_COLOR
                    p.space_after = Pt(8)
                    
                elif item['type'] == 'header':
                    if first_para:
                        p = tf.paragraphs[0]
                        first_para = False
                    else:
                        p = tf.add_paragraph()
                    p.text = item['content']
                    p.font.size = Pt(20)
                    p.font.bold = True
                    p.font.color.rgb = ACCENT_COLOR
                    p.space_before = Pt(12)
                    p.space_after = Pt(4)
                    
                elif item['type'] == 'code':
                    if first_para:
                        p = tf.paragraphs[0]
                        first_para = False
                    else:
                        p = tf.add_paragraph()
                    # Truncate long code blocks
                    code_lines = item['content'].split('\n')[:8]
                    if len(item['content'].split('\n')) > 8:
                        code_lines.append('...')
                    p.text = '\n'.join(code_lines)
                    p.font.size = Pt(12)
                    p.font.name = 'Courier New'
                    p.font.color.rgb = TEXT_COLOR
                    p.space_before = Pt(8)
                    
                elif item['type'] == 'table':
                    # Convert table to text representation
                    rows = item['content']
                    if rows:
                        # Header row
                        if first_para:
                            p = tf.paragraphs[0]
                            first_para = False
                        else:
                            p = tf.add_paragraph()
                        header = ' | '.join(rows[0])
                        p.text = header
                        p.font.size = Pt(14)
                        p.font.bold = True
                        p.font.color.rgb = ACCENT_COLOR
                        p.space_before = Pt(8)
                        
                        # Data rows
                        for row in rows[1:]:
                            p = tf.add_paragraph()
                            p.text = ' | '.join(row)
                            p.font.size = Pt(14)
                            p.font.color.rgb = TEXT_COLOR
            
            # Speaker notes
            if slide_data.speaker_notes:
                notes_slide = slide.notes_slide
                notes_tf = notes_slide.notes_text_frame
                notes_tf.text = slide_data.speaker_notes
    
    prs.save(output_path)
    print(f"‚úÖ Created presentation: {output_path}")
    print(f"   Slides: {len(slides)}")


def main():
    parser = argparse.ArgumentParser(
        description='Generate PowerPoint slides from markdown outline'
    )
    parser.add_argument(
        'input_file',
        type=Path,
        help='Path to markdown outline file'
    )
    parser.add_argument(
        '-o', '--output',
        type=Path,
        default=None,
        help='Output PPTX file path (default: same name as input with .pptx)'
    )
    
    args = parser.parse_args()
    
    if not args.input_file.exists():
        print(f"‚ùå Error: Input file not found: {args.input_file}")
        return 1
    
    # Default output path
    if args.output is None:
        args.output = args.input_file.with_suffix('.pptx')
    
    # Read markdown
    md_content = args.input_file.read_text()
    
    # Parse slides
    print(f"üìñ Parsing: {args.input_file}")
    slides = parse_markdown(md_content)
    print(f"   Found {len(slides)} slides")
    
    # Generate presentation
    print(f"üé® Generating presentation...")
    create_presentation(slides, args.output)
    
    return 0


if __name__ == '__main__':
    exit(main())
