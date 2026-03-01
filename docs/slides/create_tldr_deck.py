#!/usr/bin/env python3
"""Generate TL;DR version of AI Personas & Auto-Remediation presentation.

Creates one summary slide per section from the original 60-slide deck.
Uses the same style as onboarding-outline.pptx.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Color scheme - matching onboarding-outline.pptx
PURPLE_PRIMARY = RGBColor(0x63, 0x66, 0xF1)  # #6366F1 - section headers, titles
PURPLE_ACCENT = RGBColor(0x8B, 0x5C, 0xF6)  # #8B5CF6 - subheadings
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)  # #1F2937 - body text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_section_divider(prs, title):
    """Purple full-screen section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Purple background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PURPLE_PRIMARY
    bg.line.fill.background()

    # Title centered
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(2.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_tldr_slide(prs, section_title, original_slides, key_points):
    """TL;DR content slide with section reference."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title with slide reference
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.30), Inches(12.33), Inches(0.80)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PURPLE_PRIMARY

    # Slide reference (right-aligned, smaller)
    tb = slide.shapes.add_textbox(
        Inches(10.5), Inches(0.35), Inches(2.33), Inches(0.50)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"(Full: Slides {original_slides})"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)  # Gray
    p.alignment = PP_ALIGN.RIGHT

    # Content area
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.20), Inches(12.33), Inches(5.50)
    )
    tf = tb.text_frame
    tf.word_wrap = True

    first = True
    for item in key_points:
        if not item:
            continue

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        if item.startswith("##"):
            p.text = item[2:].strip()
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = PURPLE_ACCENT
            p.space_before = Pt(12)
        elif item.startswith("•"):
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_TEXT
        elif item.startswith("  •"):
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_TEXT
        else:
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_TEXT

    return slide


# === BUILD THE PRESENTATION ===

prs = create_presentation()

# Slide 1: Title section divider
add_section_divider(prs, "AI Personas & Auto-Remediation\nTL;DR Summary")

# Section 1: AI Foundations (Slides 1-10)
add_section_divider(prs, "AI Foundations")

add_tldr_slide(
    prs,
    "What is an AI Assistant?",
    "1-10",
    [
        "## The Foundation",
        "• LLMs are stateless - forget everything between sessions",
        '• Context window (~200K tokens) = AI\'s "working memory"',
        "• Better prompts = better results (structure, context, constraints)",
        "",
        "## The Gap",
        "• AI can explain HOW to do things, but can't DO them",
        "• MCP (Model Context Protocol) gives AI tools to take action",
        "• Prompt databases (/coffee, /start-work) provide consistency",
        "",
        "## Key Insight",
        "• Memory is the missing piece for continuity across sessions",
    ],
)

# Section 2: Why Personas (Slides 11-17)
add_section_divider(prs, "Why Personas Over Multiple Agents")

add_tldr_slide(
    prs,
    "Dynamic Personas vs Multi-Agent",
    "11-17",
    [
        "## The Problem",
        "• Multi-agent = coordination overhead, state sync, higher costs",
        "• Stateless agents create disjoint, chaotic work",
        "• Tool limit: ~80 practical, but we have 435 tools!",
        "",
        "## The Solution: Dynamic Personas",
        '• Single Claude instance wearing different "hats"',
        "• Personas are tool configuration profiles, not separate AIs",
        "• 👨‍💻 Developer (~78), 🔧 DevOps (~74), 🚨 Incident (~78), 📦 Release (~91)",
        "",
        "## How It Works",
        '• "Load devops agent" → unload current → load k8s/bonfire/quay',
        "• Cursor refreshes tool list automatically via MCP notification",
    ],
)

# Section 3: Skills (Slides 18-23)
add_section_divider(prs, "Skills - Multi-Step Workflows")

add_tldr_slide(
    prs,
    "YAML-Defined Workflows",
    "18-23",
    [
        "## What Are Skills?",
        "• YAML workflows that chain multiple tools together",
        '• "Start work" = view issue + create branch + update status',
        "• 55 production skills across 8 categories",
        "",
        "## FastMCP Skill Engine",
        "• Templating: {{ inputs.issue_key }} (Jinja2)",
        "• Conditions, compute steps, error handling",
        "• Auto-heal built into every skill",
        "",
        "## Running Skills",
        '• skill_run("start_work", \'{"issue_key": "AAP-12345"}\')',
        "• /start-work AAP-12345 (slash command)",
        "• All methods: same skill, auto-heal, memory updates",
    ],
)

# Section 4: Auto-Remediation (Slides 24-33)
add_section_divider(prs, "Auto-Remediation")

add_tldr_slide(
    prs,
    "Self-Healing Tools & Memory",
    "24-33",
    [
        "## Common Failures → Auto-Fixed",
        '• VPN disconnected → "No route to host" → vpn_connect()',
        '• Token expired → "401 Unauthorized" → kube_login()',
        "• 100% coverage: all 435 tools have @auto_heal decorators",
        "",
        "## The Learning Loop",
        "• check_known_issues() → looks up memory/learned/patterns.yaml",
        "• debug_tool() → analyzes source, proposes fix",
        "• learn_tool_fix() → saves fix forever",
        "",
        "## Key Insight",
        "• Fixes are remembered forever",
        "• The same mistake never happens twice!",
    ],
)

# Section 5: Memory (from slides 30-33 in original)
add_section_divider(prs, "Memory - Session Continuity")

add_tldr_slide(
    prs,
    "Persistent Context Across Sessions",
    "30-33",
    [
        "## Why Memory Matters",
        "• Claude is stateless by default - each session starts fresh",
        "• Memory tracks: current work, learned patterns, session history",
        "",
        "## Memory Structure",
        "• state/ - current_work.yaml, environments.yaml",
        "• learned/ - patterns.yaml, tool_fixes.yaml, runbooks.yaml",
        "• sessions/ - daily activity logs",
        "",
        "## Session Continuity",
        "• session_start() loads: active issues, today's history, patterns",
        "• check_known_issues() prevents repeating mistakes",
        "• All skills log actions to session history",
    ],
)

# Section 6: Integrations (Slides 52-56)
add_section_divider(prs, "Integrations")

add_tldr_slide(
    prs,
    "Slack Bot & IDE Extension",
    "52-56",
    [
        "## Slack Bot",
        "• Monitor channels for queries and alerts",
        "• Investigate Prometheus alerts automatically",
        "• Create Jira issues from conversations",
        "",
        "## Cursor VSCode Extension",
        "• Status bar: active issue, MR status, environment health",
        "• Tree view: Workflow Explorer in sidebar",
        "• Commands: 10+ palette commands",
        "• Notifications: toast alerts for pipeline status",
        "",
        "## Data Flow",
        "• Extension reads memory files + MCP server",
        "• D-Bus IPC for real-time Slack daemon control",
    ],
)

# Section 7: Getting Started (Slides 57-60)
add_section_divider(prs, "Getting Started")

add_tldr_slide(
    prs,
    "Quick Start & Daily Workflow",
    "57-60",
    [
        "## Installation",
        "• git clone → uv venv → uv pip install -e .",
        "• Configure Cursor MCP settings",
        '• session_start(agent="developer")',
        "",
        "## Daily Workflow",
        "• Morning: /coffee - briefing with priorities",
        "• Start: /start-work AAP-12345 - begin issue",
        "• Submit: /create-mr - code review",
        "• End: /beer - wrap up summary",
        "",
        "## Resources",
        "• README, 82 Skills, 15 Personas, 435 Tools",
        "• 130 Slash Commands documented",
    ],
)

# Save
output_path = Path(__file__).parent / "AI Personas and Auto Remediation - TLDR.pptx"
prs.save(output_path)
print(f"✅ Created: {output_path}")
print("   Title + 7 section dividers + 7 TL;DR slides = 15 slides total")
